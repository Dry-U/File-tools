#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""文档解析器模块 - 处理多种格式文档的内容提取"""
import os
import time
import logging
import PyPDF2
import pdfminer.high_level
import pdfminer.layout
try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
from docx import Document as DocxDocument
import pandas as pd
import markdown
# 尝试导入textract，如果失败则设置为None
try:
    import textract
except ImportError:
    textract = None
# from PIL import Image
# import exifread
import datetime
try:
    import win32com.client
except ImportError:
    win32com = None

class DocumentParser:
    """文档解析器类，用于提取各种格式文档的内容和元数据"""
    def _parse_pptx(self, file_path):
        """解析PPTX文件"""
        try:
            # 延迟导入以避免启动时依赖检查失败
            import pptx
            prs = pptx.Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            self.logger.warning(f"缺少 python-pptx 库，无法解析 PPTX 文件: {file_path}")
            if textract:
                try:
                    return textract.process(file_path).decode('utf-8', errors='ignore')
                except Exception as te:
                    self.logger.warning(f"无法使用textract解析PPTX文件 {file_path}: {str(te)}")
            return "错误: 缺少 python-pptx 库"
        except Exception as e:
            self.logger.error(f"PPTX解析失败 {file_path}: {str(e)}")
            # 尝试使用textract作为后备
            if textract:
                try:
                    return textract.process(file_path).decode('utf-8', errors='ignore')
                except Exception as te:
                    self.logger.warning(f"无法使用textract解析PPTX文件 {file_path}: {str(te)}")
            return f"错误: 无法解析PPTX内容\n{str(e)}"

    def __init__(self, config_loader):
        self.config_loader = config_loader
        self.logger = logging.getLogger(__name__)
        
        # 支持的文件类型映射到对应的解析函数
        self.parser_map = {
            'pdf': self._parse_pdf,
            'docx': self._parse_docx,
            'doc': self._parse_doc_win32,
            'pptx': self._parse_pptx,
            'txt': self._parse_text,
            'md': self._parse_markdown,
            'csv': self._parse_csv,
            'xlsx': self._parse_excel,
            'xls': self._parse_excel,
            # 代码文件
            'py': self._parse_text,
            'java': self._parse_text,
            'cpp': self._parse_text,
            'js': self._parse_text,
            'html': self._parse_text,
            'css': self._parse_text,
        }
    
    def _clean_text(self, text):
        """清理文本中的控制字符和乱码"""
        if not text:
            return ""
        import re
        # 1. 移除常见的不可见控制字符 (保留换行\n, 回车\r, 制表符\t)
        # \x00-\x08: NULL, SOH, STX, ETX, EOT, ENQ, ACK, BEL, BS
        # \x0b-\x0c: VT, FF
        # \x0e-\x1f: SO, SI, DLE, DC1-4, NAK, SYN, ETB, CAN, EM, SUB, ESC, FS, GS, RS, US
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
        
        # 2. 移除特殊的Unicode空白和控制字符
        # \uE000-\uF8FF: 私有使用区 (Private Use Area)，常用于图标字体或特殊符号，PDF提取常出现乱码
        # \u200b-\u200f: 零宽字符等
        # \u3000: 全角空格 (保留，中文常用)
        # \ue5d2, \ue5d3 等特定乱码
        # 扩展到所有私有使用区 (Plane 15, 16)
        # \U000F0000-\U000FFFFF: Supplementary Private Use Area-A
        # \U00100000-\U0010FFFF: Supplementary Private Use Area-B
        text = re.sub(r'[\ue000-\uf8ff\U000f0000-\U000fffff\U00100000-\U0010ffff]', '', text)
        
        # 3. 移除连续的特殊符号，如   (这些通常在私有区，已被上面规则覆盖，但为了保险)
        # 如果还有其他特定乱码，可以在这里添加
        
        # 4. 规范化空白字符：将连续的多个空格合并为一个，但保留换行结构
        text = re.sub(r'[ \t]+', ' ', text) 
        
        # 5. 合并中文之间的空格 (例如 "微 型 电 脑" -> "微型电脑")
        # 使用lookahead (?=...) 确保重叠匹配被正确处理
        # 匹配: 中文 + 空格(可能多个) + (后面紧跟中文)
        text = re.sub(r'([\u4e00-\u9fa5])\s+(?=[\u4e00-\u9fa5])', r'\1', text)
        
        # 6. 合并中文和数字/字母之间的异常空格 (针对PDF提取常见问题)
        # 例如 "论 文 1" -> "论文 1" (保留必要的间隔，但合并异常分割)
        # 这里比较激进，假设中文后紧跟空格再跟中文标点或数字通常是异常
        # 但为了安全，我们主要关注 "中文 空格 中文" 已经在上面处理了
        # 尝试处理 "中文 空格 全角字符"
        # 匹配: 中文 + 空格 + (后面紧跟全角字符)
        text = re.sub(r'([\u4e00-\u9fa5])\s+(?=[\uff00-\uffef])', r'\1', text)
        # 匹配: 全角字符 + 空格 + (后面紧跟中文)
        text = re.sub(r'([\uff00-\uffef])\s+(?=[\u4e00-\u9fa5])', r'\1', text)
        
        # 7. 移除重复的词组 (针对 "作者 简介 作者简介" 这种OCR/提取错误)
        # 匹配: 词 + 空格 + 相同的词 (仅限中文，长度2-10)
        # 例如: "作者简介 作者简介" -> "作者简介"
        text = re.sub(r'([\u4e00-\u9fa5]{2,10})\s+\1', r'\1', text)

        return text.strip()

    def extract_text(self, file_path):
        """提取文件文本内容"""
        if not os.path.exists(file_path):
            self.logger.error(f"文件不存在: {file_path}")
            return f"错误: 文件不存在"
        
        file_ext = os.path.splitext(file_path)[1].lower()[1:]  # 获取文件扩展名，去除点号
        
        # 明确拒绝图片格式，防止进入通用解析器
        if file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'svg', 'webp']:
            return ""

        try:
            text = ""
            # 尝试使用特定的解析器
            if file_ext in self.parser_map:
                text = self.parser_map[file_ext](file_path)
            else:
                # 尝试使用通用解析器作为后备
                text = self._parse_generic(file_path)
            
            return self._clean_text(text)
        except Exception as e:
            self.logger.error(f"解析文件失败 {file_path}: {str(e)}")
            return f"错误: 无法解析文件内容\n{str(e)}"
    
    def extract_metadata(self, file_path):
        """提取文件元数据"""
        if not os.path.exists(file_path):
            self.logger.error(f"文件不存在: {file_path}")
            return {"错误": "文件不存在"}
        
        metadata = {}
        
        try:
            # 获取基本文件信息
            stat_info = os.stat(file_path)
            metadata["文件路径"] = file_path
            metadata["文件名"] = os.path.basename(file_path)
            metadata["文件大小"] = f"{stat_info.st_size / 1024:.2f} KB"
            metadata["创建时间"] = datetime.datetime.fromtimestamp(stat_info.st_ctime).strftime('%Y-%m-%d %H:%M:%S')
            metadata["修改时间"] = datetime.datetime.fromtimestamp(stat_info.st_mtime).strftime('%Y-%m-%d %H:%M:%S')
            metadata["访问时间"] = datetime.datetime.fromtimestamp(stat_info.st_atime).strftime('%Y-%m-%d %H:%M:%S')
            metadata["文件类型"] = os.path.splitext(file_path)[1].upper()
            
            # 获取特定格式的元数据
            file_ext = os.path.splitext(file_path)[1].lower()[1:]
            
            if file_ext == 'pdf':
                pdf_metadata = self._extract_pdf_metadata(file_path)
                metadata.update(pdf_metadata)
            elif file_ext == 'docx':
                docx_metadata = self._extract_docx_metadata(file_path)
                metadata.update(docx_metadata)
            # 移除图片元数据提取
            # elif file_ext in ['jpg', 'jpeg', 'png', 'gif']:
            #     image_metadata = self._extract_image_metadata(file_path)
            #     metadata.update(image_metadata)
            
        except Exception as e:
            self.logger.error(f"提取元数据失败 {file_path}: {str(e)}")
            metadata["元数据提取错误"] = str(e)
        
        return metadata
    
    def _parse_pdf(self, file_path):
        """解析PDF文件"""

        # 检查文件大小，避免加载过大PDF导致内存问题
        file_size = os.path.getsize(file_path)
        max_size = 100 * 1024 * 1024  # 100MB限制
        if file_size > max_size:
            self.logger.warning(f"PDF文件过大，跳过解析 {file_path}: {file_size} bytes")
            return f"错误: PDF文件过大 ({file_size} bytes)，已跳过解析"

        text = ""
        max_text_length = 50 * 1024 * 1024  # 50MB限制输出文本

        # 1. 优先使用PyMuPDF (fitz)
        # PyMuPDF非常健壮，能处理许多pdfminer无法处理的损坏PDF，且对中文支持较好
        if fitz:
            try:
                doc = fitz.open(file_path)
                for page in doc:
                    if len(text) > max_text_length:
                        text = text[:max_text_length] + "\n... (内容已截断)"
                        break
                    # sort=True 尝试按阅读顺序排序文本块，对多栏布局有帮助
                    text += page.get_text("text", sort=True) + "\n"
                doc.close()
            except Exception as e:
                self.logger.warning(f"PyMuPDF解析PDF失败 {file_path}: {str(e)}")

        # 2. 如果PyMuPDF失败或结果为空，尝试pdfplumber (基于pdfminer.six)
        if (not text or not text.strip()) and pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        if len(text) > max_text_length:
                            text = text[:max_text_length] + "\n... (内容已截断)"
                            break
                        # extract_text() 自动处理布局
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                self.logger.warning(f"pdfplumber解析PDF失败 {file_path}: {str(e)}")

        # 3. 如果pdfplumber也失败或结果为空，尝试pdfminer.high_level
        if not text or not text.strip():
            try:
                # 尝试使用pdfminer，因为它在处理中文和布局方面通常比PyPDF2更好
                extracted_text = pdfminer.high_level.extract_text(file_path)
                if len(extracted_text) <= max_text_length:
                    text = extracted_text
                else:
                    text = extracted_text[:max_text_length] + "\n... (内容已截断)"
            except Exception as e:
                self.logger.warning(f"pdfminer解析PDF失败 {file_path}: {str(e)}")

        # 4. 如果pdfminer失败或结果为空，尝试PyPDF2
        if not text or not text.strip():
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    pages = reader.pages
                    if pages:
                        for page_num in range(len(pages)):
                            if len(text) > max_text_length:
                                text = text[:max_text_length] + "\n... (内容已截断)"
                                break
                            page = pages[page_num]
                            extracted = page.extract_text() or ""
                            if extracted:
                                text += extracted
            except Exception as e:
                self.logger.error(f"PyPDF2解析PDF失败 {file_path}: {str(e)}")

        # 5. 如果仍然为空，尝试textract
        if not text or not text.strip():
            if textract:
                try:
                    extracted = textract.process(file_path).decode('utf-8', errors='ignore')
                    if len(extracted) <= max_text_length:
                        return extracted
                    else:
                        return extracted[:max_text_length] + "\n... (内容已截断)"
                except Exception as te:
                    self.logger.warning(f"无法使用textract解析PDF文件 {file_path}: {str(te)}")
            return f"错误: 无法解析PDF内容"

        # 最终检查文本长度
        if len(text) > max_text_length:
            text = text[:max_text_length] + "\n... (内容已截断)"

        return text
    
    def _parse_docx(self, file_path):
        """解析Word文档"""
        try:
            # 检查文件大小，避免加载过大文件导致内存问题
            file_size = os.path.getsize(file_path)
            max_size = 50 * 1024 * 1024  # 50MB限制
            if file_size > max_size:
                self.logger.warning(f"Word文档过大，跳过解析 {file_path}: {file_size} bytes")
                return f"错误: Word文档过大 ({file_size} bytes)，已跳过解析"

            doc = DocxDocument(file_path)
            text = ""
            max_text_length = 10 * 1024 * 1024  # 10MB限制输出文本

            for paragraph in doc.paragraphs:
                if len(text) > max_text_length:
                    text = text[:max_text_length] + "\n... (内容已截断)"
                    break
                text += paragraph.text + '\n'

            # 最终检查文本长度
            if len(text) > max_text_length:
                text = text[:max_text_length] + "\n... (内容已截断)"

            return text
        except Exception as e:
            self.logger.error(f"DOCX解析失败 {file_path}: {str(e)}")
            # 尝试使用textract作为后备
            if textract:
                try:
                    # 检查文件大小限制
                    if file_size > max_size:
                        return f"错误: Word文档过大 ({file_size} bytes)，已跳过解析"

                    content = textract.process(file_path).decode('utf-8', errors='ignore')
                    max_content_size = 10 * 1024 * 1024  # 10MB限制内容
                    if len(content) > max_content_size:
                        content = content[:max_content_size] + "\n... (内容已截断)"
                    return content
                except Exception as te:
                    self.logger.warning(f"无法使用textract解析Word文档 {file_path}: {str(te)}")
            # 如果没有textract或解析失败，返回错误信息
            return f"错误: 无法解析Word内容\n{str(e)}"
    
    def _parse_doc_win32(self, file_path):
        """使用win32com解析.doc文件 (仅Windows)"""
        if not win32com:
            # 尝试使用textract作为后备
            if textract:
                try:
                    return textract.process(file_path).decode('utf-8', errors='ignore')
                except Exception as te:
                    self.logger.warning(f"无法使用textract解析.doc文件 {file_path}: {str(te)}")
            return "错误: 无法解析.doc内容 (缺少 pywin32 依赖或非Windows环境)"
        
        try:
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            abs_path = os.path.abspath(file_path)
            doc = word.Documents.Open(abs_path)
            text = doc.Content.Text
            doc.Close()
            # word.Quit() # 不要退出Word，可能影响其他进程，或者保持Word后台运行
            return text
        except Exception as e:
            self.logger.error(f"Win32解析.doc失败 {file_path}: {str(e)}")
            return f"错误: 无法解析.doc内容\n{str(e)}"

    def _parse_text(self, file_path):
        """解析文本文件"""
        try:
            # 检查文件大小，避免加载过大文件导致内存问题
            file_size = os.path.getsize(file_path)
            max_size = 10 * 1024 * 1024  # 10MB限制
            if file_size > max_size:
                self.logger.warning(f"文本文件过大，跳过解析 {file_path}: {file_size} bytes")
                return f"错误: 文本文件过大 ({file_size} bytes)，已跳过解析"

            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                content = file.read()
                # 限制返回内容大小以防止内存问题
                if len(content) > max_size:
                    content = content[:max_size] + "\n... (内容已截断)"
                return content
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk', errors='replace') as file:
                    content = file.read()
                    # 限制返回内容大小
                    max_size = 10 * 1024 * 1024  # 10MB限制
                    if len(content) > max_size:
                        content = content[:max_size] + "\n... (内容已截断)"
                    return content
            except:
                self.logger.error(f"文本解析失败 {file_path}: 编码错误")
                return "错误: 无法解析文本内容（编码问题）"
        except Exception as e:
            self.logger.error(f"文本解析失败 {file_path}: {str(e)}")
            return f"错误: 无法解析文本内容\n{str(e)}"
    
    def _parse_markdown(self, file_path):
        """解析Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                content = file.read()
                # 可以选择返回原始markdown或转换为HTML
                return content  # 返回原始markdown
        except Exception as e:
            self.logger.error(f"Markdown解析失败 {file_path}: {str(e)}")
            return f"错误: 无法解析Markdown内容\n{str(e)}"
    
    def _parse_csv(self, file_path):
        """解析CSV文件"""
        try:
            # 检查文件大小，避免加载过大文件导致内存问题
            file_size = os.path.getsize(file_path)
            max_size = 50 * 1024 * 1024  # 50MB限制
            if file_size > max_size:
                self.logger.warning(f"CSV文件过大，跳过完整解析 {file_path}: {file_size} bytes")
                return f"错误: CSV文件过大 ({file_size} bytes)，已跳过解析"

            df = pd.read_csv(file_path, encoding_errors='replace')

            # 限制返回内容大小
            content = df.to_string()
            if len(content) > max_size:
                content = content[:max_size] + "\n... (内容已截断)"
            return content
        except Exception as e:
            self.logger.error(f"CSV解析失败 {file_path}: {str(e)}")
            # 尝试作为普通文本文件解析
            try:
                return self._parse_text(file_path)
            except:
                return f"错误: 无法解析CSV内容\n{str(e)}"
    
    def _parse_excel(self, file_path):
        """解析Excel文件"""
        try:
            # 检查文件大小，避免加载过大文件导致内存问题
            file_size = os.path.getsize(file_path)
            max_size = 50 * 1024 * 1024  # 50MB限制
            if file_size > max_size:
                self.logger.warning(f"Excel文件过大，跳过完整解析 {file_path}: {file_size} bytes")
                return f"错误: Excel文件过大 ({file_size} bytes)，已跳过解析"

            # pandas read_excel 自动处理 .xls (需要xlrd) 和 .xlsx (需要openpyxl)
            df = pd.read_excel(file_path)

            # 限制返回内容大小
            content = df.to_string()
            max_content_size = 10 * 1024 * 1024  # 10MB限制内容
            if len(content) > max_content_size:
                content = content[:max_content_size] + "\n... (内容已截断)"
            return content
        except Exception as e:
            self.logger.error(f"Excel解析失败 {file_path}: {str(e)}")
            # 尝试使用win32com作为后备 (仅Windows)
            if win32com:
                try:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    abs_path = os.path.abspath(file_path)

                    # 检查Excel文件大小
                    file_size = os.path.getsize(file_path)
                    if file_size > 50 * 1024 * 1024:
                        excel.Quit()
                        self.logger.warning(f"Excel文件过大，跳过Win32解析 {file_path}: {file_size} bytes")
                        return f"错误: Excel文件过大 ({file_size} bytes)，已跳过解析"

                    wb = excel.Workbooks.Open(abs_path)
                    text = ""
                    for sheet in wb.Sheets:
                        if len(text) > 10 * 1024 * 1024:  # 限制输出大小
                            text += "\n... (内容已截断)"
                            break
                        try:
                            # 获取已使用区域的值
                            used_range = sheet.UsedRange
                            if used_range.Value:
                                # 将元组转换为字符串
                                text += f"Sheet: {sheet.Name}\n"
                                for row in used_range.Value:
                                    if row and len(text) <= 10 * 1024 * 1024:  # 限制输出大小
                                        text += " ".join([str(cell) for cell in row if cell is not None and cell != '']) + "\n"
                                    if len(text) > 10 * 1024 * 1024:  # 限制输出大小
                                        text += "\n... (内容已截断)"
                                        break
                        except Exception:
                            pass
                        if len(text) > 10 * 1024 * 1024:  # 限制输出大小
                            break
                    wb.Close(False)
                    excel.Quit()
                    return text
                except Exception as we:
                    self.logger.warning(f"Win32解析Excel失败 {file_path}: {str(we)}")

            # 尝试使用textract作为后备
            if textract:
                try:
                    # 检查文件大小限制
                    if file_size > 50 * 1024 * 1024:
                        return f"错误: Excel文件过大 ({file_size} bytes)，已跳过解析"

                    content = textract.process(file_path).decode('utf-8', errors='ignore')
                    max_content_size = 10 * 1024 * 1024  # 10MB限制内容
                    if len(content) > max_content_size:
                        content = content[:max_content_size] + "\n... (内容已截断)"
                    return content
                except Exception as te:
                    self.logger.warning(f"无法使用textract解析Excel文件 {file_path}: {str(te)}")
            # 如果没有textract或解析失败，返回错误信息
            return f"错误: 无法解析Excel内容\n{str(e)}"
    
    def _parse_generic(self, file_path):
        """通用解析器，用于处理不支持的文件格式"""
        file_ext = os.path.splitext(file_path)[1].lower()[1:]
        # 再次检查防止图片进入
        if file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'svg', 'webp']:
            return ""

        # 使用textract尝试提取内容
        if textract:
            try:
                return textract.process(file_path).decode('utf-8', errors='ignore')
            except Exception as te:
                self.logger.warning(f"无法使用textract解析文件 {file_path}: {str(te)}")
        # 如果没有textract或解析失败，返回错误信息
        self.logger.error(f"通用解析失败 {file_path}: 不支持的格式或缺少textract依赖")
        return f"错误: 无法解析文件内容（不支持的格式）"
    
    def _extract_pdf_metadata(self, file_path):
        """提取PDF文件特定的元数据"""
        metadata = {}
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                pdf_info = reader.metadata
                if pdf_info:
                    for key, value in pdf_info.items():
                        # 清理键名
                        clean_key = key.replace('/','')
                        metadata[f"PDF-{clean_key}"] = str(value)
                
                # 添加页数信息
                pages = reader.pages
                if pages:
                    metadata["PDF-页数"] = len(pages)
        except Exception as e:
            self.logger.error(f"提取PDF元数据失败 {file_path}: {str(e)}")
        
        return metadata
    
    def _extract_docx_metadata(self, file_path):
        """提取Word文档特定的元数据"""
        metadata = {}
        try:
            doc = DocxDocument(file_path)
            core_props = doc.core_properties
            
            if core_props.title:
                metadata["Word-标题"] = core_props.title
            if core_props.author:
                metadata["Word-作者"] = core_props.author
            if core_props.subject:
                metadata["Word-主题"] = core_props.subject
            if core_props.keywords:
                metadata["Word-关键词"] = core_props.keywords
            if core_props.comments:
                metadata["Word-注释"] = core_props.comments
            if core_props.created:
                metadata["Word-创建时间"] = core_props.created.strftime('%Y-%m-%d %H:%M:%S')
            if core_props.modified:
                metadata["Word-修改时间"] = core_props.modified.strftime('%Y-%m-%d %H:%M:%S')
                
            # 计算段落数和页数
            paragraphs = doc.paragraphs
            if paragraphs:
                metadata["Word-段落数"] = len(paragraphs)
        except Exception as e:
            self.logger.error(f"提取Word元数据失败 {file_path}: {str(e)}")
        
        return metadata
